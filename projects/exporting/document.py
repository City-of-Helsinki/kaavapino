import concurrent.futures
import threading

import datetime
import io
from html import escape
import logging
import jinja2
from jinja2 import Environment, exceptions, meta

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.parts.image import Image

from django.conf import settings
from django.http import HttpResponse
from django.utils import timezone
from docx.shared import Mm
from docxtpl import DocxTemplate, InlineImage, Listing, RichText
from PIL import Image as PImage
from ..models import Attribute, ProjectPhase, ProjectAttributeFile, ProjectPhaseSectionAttribute
from ..models.utils import create_identifier
from projects.helpers import (
    DOCUMENT_CONTENT_TYPES,
    get_file_type,
    set_kaavoitus_api_data_in_attribute_data,
    set_geoserver_data_in_attribute_data,
    set_ad_data_in_attribute_data,
    set_automatic_attributes,
)
from projects.models import ProjectDocumentDownloadLog

log = logging.getLogger(__name__)

MAX_WIDTH_MM = 170  # Max InlineImage width
DEFAULT_IMG_DPI = (72, 72)  # For cases where dpi value is not available in metadata


def _get_raw_value(value, attribute):
    if attribute.value_type == Attribute.TYPE_DATE and isinstance(value, str):
        return datetime.datetime.strptime(value, "%Y-%m-%d").date()
    else:
        return value


def _set_fieldset_path(fieldset_path, attribute_data_display, identifier, value):
    parent = fieldset_path[0]["parent"].identifier
    index = fieldset_path[0]["index"]

    for idx, fieldset_item in enumerate(attribute_data_display[parent]):
        if fieldset_item["index"] == index:
            fieldset_item[identifier] = value
            attribute_data_display[parent][idx] = fieldset_item


def get_top_level_attribute(attribute):
    if not attribute.fieldsets.count():
        return attribute
    else:
        return get_top_level_attribute(attribute.fieldsets.first())


def get_attribute_subtitle(target_identifier, target_phase_id, project):
    try:
        if target_identifier == "vastuuhenkilo_nimi":
            target_identifier = "vastuuhenkilo_nimi_readonly"
        projectphasesectionattribute = ProjectPhaseSectionAttribute.objects.filter(
            attribute__identifier=target_identifier,
            section__phase=target_phase_id
        )
        #Could have multiple sections but we need the closest one to the current phase
        if len(projectphasesectionattribute) > 0:
            section = projectphasesectionattribute.filter(index__gte=project.phase.index).first()
            return section.section.name
        else:
            return projectphasesectionattribute.section.name
    except ProjectPhaseSectionAttribute.DoesNotExist:
        return None


def get_closest_phase(project, identifier, parent_identifier=None):
    if parent_identifier == "vastuuhenkilo_nimi":
        parent_identifier = "vastuuhenkilo_nimi_readonly"
    phases = ProjectPhase.objects.filter(
        sections__attributes__identifier=identifier,
        project_subtype=project.subtype,
    ).order_by("index")
    #If not found in the current identifier, check the parent
    if not phases and parent_identifier is not None:
        phases = ProjectPhase.objects.filter(
            sections__attributes__identifier=parent_identifier,
            project_subtype=project.subtype,
        ).order_by("index")
    # Returning the closest open phase if found,
    # otherwise return the last phase when the attribute
    # was editable
    phase = phases.filter(index__gte=project.phase.index).first()
    return phase or phases.reverse().first()


def get_rich_text_display_value(value, preview=False, **text_args):
    if not value:
        return RichText("Tieto puuttuu", **text_args)

    rich_text = RichText(None)

    try:
        url_id = text_args.get("url_id", None)
        color = text_args.get("color", None)

        operations = value["ops"]

        for index, operation in enumerate(operations, start=1):
            insert = operation.get("insert", None)

            if index == len(operations):
                insert = insert.rstrip() if insert else None

            if not insert:
                continue

            attributes = operation.get("attributes", None)
            if not attributes:
                rich_text.add(insert,
                              url_id=url_id,
                              color=color
                              )
                continue
            _color = get_color(preview,color,attributes)
            _size = attributes.get("size", None)
            _script = attributes.get("script", None)
            _sub = get_sub(_script)
            _super = get_super(_script)
            _bold = attributes.get("bold", False)
            _italic = attributes.get("italic", False)
            _underline = attributes.get("underline", False)
            _strike = attributes.get("strike", False)
            _font = attributes.get("font", None)

            rich_text.add(insert,
                          color=_color,
                          size=_size,
                          subscript=_sub,
                          superscript=_super,
                          bold=_bold,
                          italic=_italic,
                          underline=_underline,
                          strike=_strike,
                          font=_font,
                          url_id=url_id
                          )
    except Exception as exc:
        log.error("Error while formatting RichText value", exc)
        return RichText("Virhe arvoa generoitaessa", **text_args)

    return rich_text

def get_color(preview,color,attributes):
    if not preview:
        return None
    elif color:
        return color
    else:
        return attributes.get("color", None)
def get_sub(_script):
    if _script == "sub":
        return True
    else:
        return False
def get_super(_script):
    if _script == "super":
        return True
    else:
        return False
def render_template(project, document_template, preview):
    doc_type = get_file_type(document_template.file.path)

    if doc_type == 'docx':
        doc = DocxTemplate(document_template.file)
    else:
        doc = None

    attribute_data_display = {}
    attribute_element_data = {}
    attributes = {a.identifier: a for a in Attribute.objects.all()}

    def get_display_and_raw_value(attribute, value, ignore_multiple_choice=False):
        empty = False
        text_args = None
        element_data = {}

        if attribute.value_type in [Attribute.TYPE_FIELDSET, Attribute.TYPE_INFO_FIELDSET]:
            result = []
            for index, fieldset_item in enumerate(value) if value else []:
                fieldset_object = {}
                for k, v in fieldset_item.items():
                    item_attr = attributes.get(k)
                    if fieldset_item.get("_deleted") or not item_attr:
                        continue

                    display_value, raw_value, element_data = \
                        get_display_and_raw_value(item_attr, v)

                    fieldset_object[k] = display_value

                    if item_attr.value_type not in [Attribute.TYPE_FIELDSET, Attribute.TYPE_INFO_FIELDSET]:
                        fieldset_object[f"{k}__raw"] = raw_value

                if fieldset_object:
                    fieldset_object["index"] = index
                    result.append(fieldset_object)

            return (result, value, element_data)
        elif attribute.multiple_choice and not ignore_multiple_choice:
            value = value or []
            display_list = [
                get_display_and_raw_value(
                    attribute, i, ignore_multiple_choice=True
                )[0] for i in value
            ]
            raw_list = [
                _get_raw_value(i, attribute) for i in value
            ]
            all_element_data = [
                get_display_and_raw_value(
                    attribute, i, ignore_multiple_choice=True
                )[2] for i in value
            ]
            if all_element_data:
                element_data = all_element_data[0]

            return (display_list, raw_list, element_data)

        if attribute.value_type == Attribute.TYPE_IMAGE and value:
            if doc_type == 'docx':
                try:
                    with PImage.open(value) as img:
                        width_px = img.width
                        dpi = float(img.info.get('dpi', DEFAULT_IMG_DPI)[0])
                        dpi = dpi if dpi > 0 else DEFAULT_IMG_DPI[0]
                    width_mm = int((width_px/dpi) * 25.4)
                    display_value = InlineImage(doc, value, width=Mm(MAX_WIDTH_MM) if width_mm > MAX_WIDTH_MM else Mm(width_mm))
                except FileNotFoundError:
                    log.error(f'Image not found at {value}')
                    display_value = None
            else:
                display_value = value
        else:
            display_value = attribute.get_attribute_display(value)

        if display_value is None or display_value == "":
            empty = True
            if preview:
                display_value = "Tieto puuttuu"

        if attribute.value_type == Attribute.TYPE_LONG_STRING:
            display_value = Listing(display_value)
        if attribute.value_type == Attribute.TYPE_IMAGE:
            pass
        else:
            if preview:
                target_property = None

                if attribute.static_property:
                    target_identifier = None
                    target_property = attribute.static_property

                target_identifier = \
                    get_top_level_attribute(attribute).identifier

                target_phase_id = None
                target_section_name = None
                if target_identifier:
                    try:
                        target_phase_id = get_closest_phase(project, attribute.identifier,target_identifier).id
                        target_section_name = get_attribute_subtitle(target_identifier, target_phase_id, project)
                    except AttributeError:
                        pass

                # attribute: editable attribute field
                # phase: closest phase where attribute field is located
                # property: editable project model field
                # view: default | deadlines | floorarea
                edit_url = settings.DOCUMENT_EDIT_URL_FORMAT.replace(
                    "<pk>", str(project.pk),
                )

                if attribute.projectphasedeadlinesectionattribute_set.count():
                    view = "deadlines"
                elif attribute.projectfloorareasectionattribute_set.count():
                    view = "floorarea"
                else:
                    view = "default"

                edit_url += f"?view={view}"
                edit_url += f"&attribute={target_identifier}" if target_identifier else ""
                edit_url += f"&phase={target_phase_id}" if target_phase_id else ""
                edit_url += f"&section={target_section_name}" if target_section_name else ""
                edit_url += f"&property={target_property}" if (target_property and not target_identifier) else ""

                with build_url_id_lock:
                    text_args = {
                        "color": "#d0c873" if empty else "#79a6b5",
                        "url_id": doc.build_url_id(edit_url) if doc_type == 'docx' else edit_url,
                    }
            else:
                text_args = {}

            if doc_type == 'docx':
                if attribute.value_type in [Attribute.TYPE_RICH_TEXT, Attribute.TYPE_RICH_TEXT_SHORT]:
                    display_value = get_rich_text_display_value(value, preview, **text_args)
                else:
                    display_value = RichText(display_value, **text_args)

        return (display_value, _get_raw_value(value, attribute), text_args)

    attribute_data = project.attribute_data
    try:
        set_kaavoitus_api_data_in_attribute_data(attribute_data)
    except Exception:
        pass

    set_geoserver_data_in_attribute_data(attribute_data)
    set_ad_data_in_attribute_data(attribute_data)
    set_automatic_attributes(attribute_data)

    build_url_id_lock = threading.Lock()

    def process_attribute(attr):
        value = attribute_data.get(attr.identifier)
        identifier = attr.identifier
        display_value, raw_value, element_data = get_display_and_raw_value(attr, value)
        return identifier, display_value, raw_value, element_data

    full_attribute_data = [(attr, attribute_data.get(attr.identifier)) for attr in Attribute.objects.all()]

    with concurrent.futures.ThreadPoolExecutor() as executor:
        future_to_attr = {executor.submit(process_attribute, attr): attr for attr, _ in full_attribute_data}
        for future in concurrent.futures.as_completed(future_to_attr):
            attr = future_to_attr[future]

            identifier, display_value, raw_value, element_data = future.result()
            attribute_data_display[identifier] = display_value
            attribute_element_data[identifier] = element_data
            if attr.value_type != Attribute.TYPE_FIELDSET:
                attribute_data_display[identifier + "__raw"] = raw_value

    attribute_files = ProjectAttributeFile.objects \
        .filter(project=project, archived_at=None) \
        .order_by(
            "fieldset_path_str",
            "attribute__pk",
            "project__pk",
            "-created_at",
        ) \
        .distinct("fieldset_path_str", "attribute__pk", "project__pk")

    for attribute_file in attribute_files:
        # only image formats supported by docx/pptx can be used
        image_formats = [
            "bmp",
            "emf",
            "emz",
            "eps",
            "fpix", "fpx",
            "gif",
            "jpg", "jpeg", "jfif", "jpeg-2000",
            "pict", "pct",
            "png",
            "pntg",
            "psd",
            "qtif",
            "sgi",
            "tga", "tpic",
            "tiff", "tif",
            "wmf",
            "wmz",
        ]

        file_format_is_supported = \
            attribute_file.file.path.split('.')[-1].lower() in image_formats

        if file_format_is_supported:
            display_value, __, __ = get_display_and_raw_value(
                attribute_file.attribute,
                attribute_file.file.path,
            )
        elif preview:
            display_value = "Kuvan tiedostotyyppiä ei tueta"
        else:
            continue

        if not attribute_file.fieldset_path:
            attribute_data_display[
                attribute_file.attribute.identifier
            ] = display_value
        else:
            _set_fieldset_path(
                attribute_file.fieldset_path,
                attribute_data_display,
                attribute_file.attribute.identifier,
                display_value,
            )

    # Add preview information to attribute_data_display so that it can be used as condition in documents
    attribute_data_display.update({'is_preview': preview})

    jinja_env = jinja2.Environment()
    jinja_env.filters['distinct'] = distinct

    output = None

    if doc_type == 'docx':
        try:
            doc.render(attribute_data_display, jinja_env)
            output = io.BytesIO()
            doc.save(output)
        except Exception as exc:
            log.error('Error while rendering document', exc)
    else:
        data = {
            'data': attribute_data_display,
            'element_data': attribute_element_data,
        }
        doc = PptxTemplate(document_template.file, data, env=jinja_env)
        output = doc.save()

    if not preview:
        ProjectDocumentDownloadLog.objects.create(
            project=project,
            document_template=document_template,
            phase=project.phase.common_project_phase,
        )

    return output.getvalue() if output else "error"


# Custom filter for filtering objects from list by unique key
def distinct(value, key):
    if value and type(value) is list:
        checked = set()
        filtered = [e for e in value
            if e.get(key) and e.get(key) not in checked
            and not checked.add(e.get(key))
        ]

        return filtered

    return value


def get_document_response(project, document_template, filename=None):
    if filename is None:
        filename = "{}-{}-{}".format(
            create_identifier(project.name),
            document_template.name,
            timezone.now().date(),
        )

    doc_type = get_file_type(document_template.file.path)
    output = render_template(project, document_template)
    response = HttpResponse(
        output,
        content_type=DOCUMENT_CONTENT_TYPES[doc_type]
    )
    response["Content-Disposition"] = "attachment; filename={}.{}".format(filename, doc_type)
    response["Access-Control-Allow-Origin"] = "*"
    return response


class PptxTemplate:
    def __init__(self, template_file, data, env):
        self.template_file = template_file
        self.data = data['data']
        self.element_data = data['element_data']
        self.env = env

    def save(self):
        doc = Presentation(self.template_file)
        for slide in doc.slides:
            self.render_slide(slide)

        output = io.BytesIO()
        doc.save(output)
        return output

    def render_slide(self, slide):
        for shape in slide.shapes:
            self.render_shape(shape)

    def render_shape(self, shape):
        if shape.has_text_frame:
            self.render_text_frame(shape.text_frame)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self.render_picture(shape)

    def render_picture(self, shape):
        img = self.data.get(shape.name)
        if img:
            parent = shape._parent
            try:
                new_picture = parent.add_picture(
                    img,
                    shape.left,
                    shape.top,
                    height=shape.height
                )
            except FileNotFoundError:
                pass
            parent.element.remove(shape._element)

    def render_text_frame(self, text_frame):
        for paragraph in text_frame.paragraphs:
            self.render_paragraph(paragraph)

    def render_paragraph(self, paragraph):
        original_text = paragraph.text
        env = self.env.from_string(original_text)
        rendered = env.render(self.data)

        # "Clear" all but one run and use it to retain template's text styles
        p = paragraph._p
        for idx, run in enumerate(paragraph.runs):
            if idx == 0 and rendered.strip():
                continue
            p.remove(run._r)

        if paragraph.runs:
            template = self.env.parse(original_text)
            vars = self.get_cleaned_context_variables(template)
            rendered = self.clean_empty_lines(rendered)
            paragraph.runs[0].text = rendered

            if len(vars) == 1:
                link_data = self.element_data.get(vars.pop())
                if link_data:
                    font = paragraph.runs[0].font
                    paragraph.runs[0].hyperlink.address = link_data['url_id']

        # Clean up potential empty paragraph
        if not paragraph.text.strip():
            p = paragraph._p
            p.getparent().remove(p)

    def clean_empty_lines(self, rendered_text):
        return '\n'.join([txt for txt in rendered_text.split('\n') if txt.strip()])

    def get_cleaned_context_variables(self, template):
        return set([
            v for v in meta.find_undeclared_variables(template)
            if '__raw' not in v and 'sorted' != v
        ])
